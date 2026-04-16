"""
Gerador de Estudo Técnico — Template-based 100% dinâmico.

Recebe um ResultadoHistorico (já calculado pelo AnalisadorHistorico)
e preenche o template PPTX com os dados do cliente.
Responsabilidade única: formatação e geração do documento.
Nenhuma regra de negócio ou cálculo aqui.
"""
from __future__ import annotations
import io, os, logging, subprocess, tempfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches

from src.parsers.analyzer_historico import ResultadoHistorico

logger = logging.getLogger(__name__)


def _f(v, d=2):
    try:
        return f"{float(v):,.{d}f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except Exception:
        return "—"

def _fi(v):
    return _f(v, 0)

def _assets():
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

def _replace_in_shape(shape, old, new):
    if not shape.has_text_frame: return False
    if old not in shape.text_frame.text: return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
        if para.runs:
            combined = "".join(r.text for r in para.runs)
            if old in combined:
                para.runs[0].text = combined.replace(old, new)
                for r in para.runs[1:]: r.text = ""
                return True
    return False

def _replace_all(prs, old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, old, new)
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if old in cell.text:
                            for p in cell.text_frame.paragraphs:
                                for r in p.runs:
                                    if old in r.text:
                                        r.text = r.text.replace(old, new)

def _set_autofit(shape):
    """Habilita auto-fit de texto na shape (encolhe para caber no box)."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            # Remove normAutofit ou spAutoFit existentes
            for child in list(bodyPr):
                if child.tag in (qn('a:normAutofit'), qn('a:spAutofit'), qn('a:noAutofit')):
                    bodyPr.remove(child)
            # Adiciona normAutofit (encolhe fonte para caber)
            bodyPr.append(etree.fromstring('<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
    except Exception:
        pass


def _set_cell(tbl, row, col, text):
    cell = tbl.cell(row, col)
    for p in cell.text_frame.paragraphs:
        if p.runs:
            p.runs[0].text = str(text)
            for r in p.runs[1:]: r.text = ""
            return
    cell.text = str(text)

def _delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


# ═══════════════════════════════════════════════════════════════════════════════

def gerar_estudo_pdf(
    resultado: ResultadoHistorico,
    valor_mensal: int = 500,
    comissao: int = 30,
    pdf_screenshot_bytes: bytes = None,
) -> io.BytesIO:
    """
    Preenche o template PPTX com dados do ResultadoHistorico e converte para PDF.
    """
    r = resultado
    meses_pt = {1:"Janeiro",2:"Fevereiro",3:"Março",4:"Abril",5:"Maio",6:"Junho",
                7:"Julho",8:"Agosto",9:"Setembro",10:"Outubro",11:"Novembro",12:"Dezembro"}
    mes_atual = meses_pt.get(datetime.now().month, "")

    tmpl = os.path.join(_assets(), "template.pptx")
    if not os.path.exists(tmpl):
        raise FileNotFoundError(f"Template não encontrado: {tmpl}")
    prs = Presentation(tmpl)

    # SLIDE 1: CAPA
    _replace_all(prs, "COMETAIS INDUSTRIA E COMERCIO DE METAIS LTDA", r.nome[:60])
    _replace_all(prs, "02.896.727/0003-96", r.cnpj or "")
    _replace_all(prs, "0502485-4", r.uc)

    # SLIDE 3: SCREENSHOT DA FATURA
    try:
        s3 = list(prs.slides)[2]
        for shape in list(s3.shapes):
            if shape.shape_type == 13 and shape.width > Inches(3):
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                if pdf_screenshot_bytes:
                    s3.shapes.add_picture(io.BytesIO(pdf_screenshot_bytes), left, top, width, height)
                    logger.info("[Estudo] Screenshot embutido")
                break
    except Exception as e:
        logger.warning(f"[Estudo] Erro slide 3: {e}")

    # SLIDE 3: PONTOS DE ATENÇÃO
    pts = r.pontos_atencao
    _replace_all(prs, "Demanda contratada atual.", pts[0] if len(pts) > 0 else "")
    _replace_all(prs, "Consumo na Ponta (Entre as 20h às 23h) muito elevado, podendo ser feito um estudo de BESS",
                 pts[1] if len(pts) > 1 else "")
    _replace_all(prs, "(Armazenamento de energia em baterias).", "")
    _replace_all(prs, "Demanda não utilizada recorrente, muito desperdício todos os meses. Podendo ser ajustada com sazonalidade.",
                 pts[2] if len(pts) > 2 else "")
    _replace_all(prs, "Energia reativa, multa devido ao baixo fator  de potência, corrigir com Banco de Capacitores e estudar o Filtro Capacitivo para protege-lo e os demais equipamentos da indústria.",
                 pts[3] if len(pts) > 3 else "")
    _replace_all(prs, "Multas por atrasos de pagamentos constantes, requer uma gestão ativa nas contas.",
                 pts[4] if len(pts) > 4 else "")
    _replace_all(prs, "Descrição de Gradeza são os registros feitos pelo medidor da UC, aqui acontecem as auditorias retroativas.",
                 pts[5] if len(pts) > 5 else "")

    # SLIDE 5: REMOVER FREEFORM OVAL decorativa (shape index 4)
    try:
        slide5 = list(prs.slides)[4]
        shapes_to_remove = [
            s for s in slide5.shapes
            if s.shape_type == 5 and not s.has_text_frame  # FREEFORM sem texto
            or (s.shape_type == 5 and s.has_text_frame and not s.text_frame.text.strip())
        ]
        for s in shapes_to_remove:
            sp = s._element
            sp.getparent().remove(sp)
        if shapes_to_remove:
            logger.info(f"[Estudo] {len(shapes_to_remove)} freeform(s) removida(s) do slide 5")
    except Exception as e:
        logger.warning(f"[Estudo] Erro remover freeform slide 5: {e}")

    # SLIDE 4: REMOVER (específico Cometais)
    try:
        _delete_slide(prs, 3)
        logger.info("[Estudo] Slide 4 removido")
    except Exception as e:
        logger.warning(f"[Estudo] Erro remover slide 4: {e}")

    # SLIDE 5→4: RESUMO TÉCNICO
    _replace_all(prs, "CONSOLIDADO DOS ÚLTIMOS 12 MESES:", f"CONSOLIDADO DOS ÚLTIMOS {r.n_faturas} MESES:")

    for sl in prs.slides:
        for sh in sl.shapes:
            if hasattr(sh, "has_table") and sh.has_table and len(sh.table.rows) >= 8:
                tbl = sh.table
                tusd_pct, te_pct = 63.43, 36.57
                tusd_v = r.custo_anual * tusd_pct / 100
                te_v   = r.custo_anual * te_pct / 100
                _set_cell(tbl,0,0,"Custo do Fio (TUSD — AME):")
                _set_cell(tbl,0,1,f"R$ {_f(tusd_v)}")
                _set_cell(tbl,0,2,f"{tusd_pct:.2f}%")
                _set_cell(tbl,1,0,"Custo da Energia (TE — SAFIRA):")
                _set_cell(tbl,1,1,f"R$ {_f(te_v)}")
                _set_cell(tbl,1,2,f"{te_pct:.2f}%")
                _set_cell(tbl,2,0,"")
                _set_cell(tbl,2,1,f"R$ {_f(r.custo_anual)}")
                _set_cell(tbl,2,2,"100%")
                _set_cell(tbl,3,0,f"DESPERDÍCIO DOS ÚLTIMOS {r.n_faturas} MESES:")
                _set_cell(tbl,3,1,""); _set_cell(tbl,3,2,"")

                pot = r.potencial_periodo
                def _pct(v): return round(v/pot*100,2) if pot>0 else 0

                desp = []
                if r.demanda_ociosa_r > 0:
                    desp.append(("Demanda não utilizada:", f"R$ {_f(r.demanda_ociosa_r)}", f"{_pct(r.demanda_ociosa_r):.2f}%"))
                if r.demanda_ultrapassagem_r > 0:
                    desp.append(("Demanda ultrapassada:", f"R$ {_f(r.demanda_ultrapassagem_r)}", f"{_pct(r.demanda_ultrapassagem_r):.2f}%"))
                if r.reativo_r > 0:
                    desp.append(("Energia reativa:", f"R$ {_f(r.reativo_r)}", f"{_pct(r.reativo_r):.2f}%"))
                if r.multas_atraso_r > 0:
                    desp.append(("Multas por atrasos:", f"R$ {_f(r.multas_atraso_r)}", f"{_pct(r.multas_atraso_r):.2f}%"))
                if r.icms_recuperavel_r > 0:
                    desp.append(("ICMS recuperável:", f"R$ {_f(r.icms_recuperavel_r)}", f"{_pct(r.icms_recuperavel_r):.2f}%"))

                for i in range(6):
                    row = 4 + i
                    if row >= len(tbl.rows): break
                    if i < len(desp):
                        _set_cell(tbl,row,0,desp[i][0]); _set_cell(tbl,row,1,desp[i][1]); _set_cell(tbl,row,2,desp[i][2])
                    elif i == len(desp):
                        _set_cell(tbl,row,0,""); _set_cell(tbl,row,1,f"R$ {_f(r.potencial_anual)}"); _set_cell(tbl,row,2,"100%")
                    else:
                        _set_cell(tbl,row,0,""); _set_cell(tbl,row,1,""); _set_cell(tbl,row,2,"")
                break

    # SLIDE 6→5: COM vs SEM GESTÃO
    _replace_all(prs, "-R$ 183.602,66", f"-R$ {_f(r.potencial_anual)}")
    _replace_all(prs, "+R$ 183.602,66", f"+R$ {_f(r.potencial_anual)}")

    # SLIDE 7→6: CRONOGRAMA
    _replace_all(prs, "Cronograma de Ações (Abril)", f"Cronograma de Ações ({mes_atual})")
    acoes = r.acoes

    _replace_all(prs, "Buscar pagamentos indevidos e pedir devolução em dobro",
                 acoes[0][1] if len(acoes)>0 else "")
    _replace_all(prs, "Portabilidade  para Âmbar Energia no Mercado Livre",
                 acoes[2][0] if len(acoes)>2 else "Migração\npara o\nMercado Livre")
    _replace_all(prs, "Atual proprietária da AME. Reduzir preço",
                 acoes[2][1].split("\n")[0] if len(acoes)>2 else "")
    _replace_all(prs, "e estreitar o relacionamento",
                 "\n".join(acoes[2][1].split("\n")[1:]) if len(acoes)>2 else "")
    _replace_all(prs, "Estudar  o consumo no",
                 "Ajustar a\ndemanda" if r.utilizacao_demanda<85 else "Estudar\no consumo no")
    _replace_all(prs, "horário de ponta",
                 "contratada\nociosa" if r.utilizacao_demanda<85 else "horário de\nponta")
    _replace_all(prs, "Instalar analisador, levantar a carga e fazer um estudo do BESS",
                 acoes[1][1] if len(acoes)>1 else "")
    _replace_all(prs, "Ajustar Banco de Capacitores existentes",
                 acoes[3][0] if len(acoes)>3 else "")
    _replace_all(prs, "Corrigir horário indutivo dos bancos URGENTE",
                 acoes[3][1] if len(acoes)>3 else "")
    _replace_all(prs, "Fazer laudo para separar o que é produção",
                 acoes[4][1] if len(acoes)>4 else "")

    cosip_m = r.cosip_media
    _replace_all(prs, "Instalar Filtro Capacitivo",
                 "Contestar\nCOSIP" if cosip_m>300 else "Relatório\nMensal de")
    _replace_all(prs, "e corrigir Banco de Capacitores",
                 "junto à\nPrefeitura" if cosip_m>300 else "Resultado")
    _replace_all(prs, "Para proteção dos Bancos de",
                 "Valor médio" if cosip_m>300 else "Prestação de contas")
    _replace_all(prs, "Capacitores e motores",
                 f"R$ {_fi(cosip_m)}/mês" if cosip_m>300 else "com economia e metas")
    _replace_all(prs, "Se não houver aumento do consumo, ajustar",
                 acoes[-1][1] if acoes else "")
    _replace_all(prs, "Ajustar a demanda contratada ociosa",
                 acoes[-1][0] if acoes else "")

    # Auto-fit em todas as shapes de texto do cronograma
    try:
        # Cronograma é slide 7→6 (índice 5 após deleção do slide 4)
        crono_slide = list(prs.slides)[5]
        for shape in crono_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                _set_autofit(shape)
    except Exception as e:
        logger.warning(f"[Estudo] Auto-fit cronograma: {e}")

    # SLIDE 9→8: MONITORAMENTO
    _replace_all(prs, "João Gomes", r.nome[:25])

    # SLIDE 10→9: PROPOSTA
    _replace_all(prs, "R$ 1.500,00 nos três primeiros meses", f"R$ {_fi(valor_mensal)},00/mês por UC")
    _replace_all(prs, "R$ 2.500,00 no quarto mês em diante", f"{comissao}% dos valores recuperados")
    _replace_all(prs, "30% dos valores que forem recuperados", "")

    # CONVERTER PPTX → PDF
    with tempfile.TemporaryDirectory() as tmp:
        px = os.path.join(tmp, "estudo.pptx")
        prs.save(px)
        logger.info(f"[Estudo] PPTX: {os.path.getsize(px)} bytes, {len(list(prs.slides))} slides")
        profile = os.path.join(tmp, "lo_profile")
        os.makedirs(profile, exist_ok=True)
        try:
            subprocess.run(
                ["soffice","--headless","--norestore","--nofirststartwizard",
                 f"-env:UserInstallation=file://{profile}",
                 "--convert-to","pdf","--outdir",tmp,px],
                check=True, timeout=60, capture_output=True
            )
            pdfs = [f for f in os.listdir(tmp) if f.endswith(".pdf")]
            if pdfs:
                with open(os.path.join(tmp, pdfs[0]), "rb") as f:
                    buf = io.BytesIO(f.read())
                buf.seek(0)
                logger.info(f"[Estudo] PDF: {buf.getbuffer().nbytes} bytes, UC {r.uc}")
                return buf
        except Exception as e:
            logger.warning(f"[Estudo] soffice falhou ({e}), retornando PPTX")
        with open(px, "rb") as f:
            buf = io.BytesIO(f.read())
        buf.seek(0)
        return buf


gerar_estudo_pptx = gerar_estudo_pdf
